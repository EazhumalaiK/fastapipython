from fastapi import FastAPI, UploadFile, File, Form, HTTPException
from fastapi.responses import JSONResponse, FileResponse
from fastapi.middleware.cors import CORSMiddleware
from fastapi.staticfiles import StaticFiles
import os
import shutil
from pptx import Presentation
from PIL import Image, ImageDraw, ImageFont
from io import BytesIO
import traceback
import logging

app = FastAPI()

# CORS configuration
app.add_middleware(
    CORSMiddleware,
    allow_origins=["https://pptreview.netlify.app"],  # Allow only the Netlify origin
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

OUTPUT_DIR = "slides"
ORIGINAL_DIR = "original_slides"

# Store comments per slide in memory (replace with DB in production)
comments_store = {}

def draw_slide_content(slide, width, height, output_path):
    img = Image.new("RGB", (width, height), color="white")
    draw = ImageDraw.Draw(img)
    font = ImageFont.load_default()

    for shape in slide.shapes:
        if shape.shape_type == 13:  # Picture (check for image type)
            image = shape.image
            try:
                image_bytes = image.blob
                pil_img = Image.open(BytesIO(image_bytes))
            except IOError as e:
                logger.warning(f"Error opening image in shape: {str(e)}. Skipping this image.")
                continue  # Skip problematic image

            left = shape.left // 9525
            top = shape.top // 9525
            shape_width = shape.width // 9525
            shape_height = shape.height // 9525

            pil_img = pil_img.resize((shape_width, shape_height))
            img.paste(pil_img, (left, top))

        elif shape.has_text_frame:
            text_frame = shape.text_frame
            text = ""
            for paragraph in text_frame.paragraphs:
                para_text = "".join(run.text for run in paragraph.runs)
                if para_text.strip():
                    text += para_text + "\n"

            if not text.strip():
                continue

            left = shape.left // 9525
            top = shape.top // 9525

            draw.text((left, top), text, fill="black", font=font)

    img.save(output_path)


def overlay_comments(slide_image_path, comments):
    img = Image.open(slide_image_path).convert("RGBA")
    overlay = Image.new("RGBA", img.size, (255, 255, 255, 0))
    draw = ImageDraw.Draw(overlay)
    font = ImageFont.load_default()

    y_offset = 10
    for comment in comments:
        draw.text((10, y_offset), comment, fill=(255, 0, 0, 255), font=font)  # red text
        y_offset += 20

    combined = Image.alpha_composite(img, overlay).convert("RGB")
    return combined


@app.post("/convert-ppt")
async def convert_ppt(file: UploadFile = File(...)):
    try:
        if not file.filename.lower().endswith(".pptx"):
            return JSONResponse(content={"error": "Only .pptx files are supported."}, status_code=400)

        # Save uploaded file temporarily
        file_path = f"temp_{file.filename}"
        with open(file_path, "wb") as f:
            f.write(await file.read())

        # Clear previous output dirs
        for d in [OUTPUT_DIR, ORIGINAL_DIR]:
            if os.path.exists(d):
                shutil.rmtree(d)
            os.makedirs(d)

        prs = Presentation(file_path)
        slide_width = prs.slide_width // 9525
        slide_height = prs.slide_height // 9525

        for idx, slide in enumerate(prs.slides, start=1):
            orig_path = os.path.join(ORIGINAL_DIR, f"slide_{idx}.png")
            draw_slide_content(slide, slide_width, slide_height, orig_path)

            # Initially, copy original image to output dir (no comments)
            shutil.copy(orig_path, os.path.join(OUTPUT_DIR, f"slide_{idx}.png"))

        os.remove(file_path)

        # Clear comments store
        comments_store.clear()

        slide_urls = [f"/slides/slide_{i+1}.png" for i in range(len(prs.slides))]
        return {"slideCount": len(prs.slides), "imageUrls": slide_urls}

    except Exception as e:
        # Log the error in detail
        logger.error(f"Error in convert_ppt: {str(e)}")
        traceback_str = traceback.format_exc()
        raise HTTPException(status_code=500, detail=f"Internal Server Error: {str(e)}\n\n{traceback_str}")


@app.post("/slides/{slide_number}/comment")
async def add_comment(slide_number: int, comment: str = Form(...)):
    try:
        orig_image_path = os.path.join(ORIGINAL_DIR, f"slide_{slide_number}.png")
        if not os.path.exists(orig_image_path):
            raise HTTPException(status_code=404, detail="Slide not found")

        # Add comment to store
        comments_store.setdefault(slide_number, []).append(comment)

        # Overlay all comments on original image
        updated_img = overlay_comments(orig_image_path, comments_store[slide_number])

        # Save updated image in OUTPUT_DIR
        updated_path = os.path.join(OUTPUT_DIR, f"slide_{slide_number}.png")
        updated_img.save(updated_path)

        return {"message": "Comment added", "updatedImageUrl": f"/slides/slide_{slide_number}.png"}

    except Exception as e:
        # Log the error in detail
        logger.error(f"Error in add_comment: {str(e)}")
        traceback_str = traceback.format_exc()
        raise HTTPException(status_code=500, detail=f"Internal Server Error: {str(e)}\n\n{traceback_str}")


@app.get("/slides/{slide_number}/comments")
async def get_comments(slide_number: int):
    return {"comments": comments_store.get(slide_number, [])}


# Serve images from the OUTPUT_DIR folder
app.mount("/slides", StaticFiles(directory=OUTPUT_DIR), name="slides")
